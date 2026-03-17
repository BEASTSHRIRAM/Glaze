"""File content extraction and processing."""
import logging
import io
import re
from typing import List
try:
    from PyPDF2 import PdfReader
except ImportError:
    from pypdf import PdfReader
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)


class FileProcessor:
    """Handles file content extraction and text chunking."""
    
    def extract_text(self, file_content: bytes, mime_type: str) -> str:
        """
        Extract text content from various file formats.
        
        Args:
            file_content: File content as bytes
            mime_type: MIME type of the file
            
        Returns:
            Extracted text content
            
        Raises:
            Exception: If extraction fails
        """
        try:
            if mime_type == "application/pdf":
                return self._extract_pdf(file_content)
            elif mime_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                               "application/vnd.google-apps.document"]:
                return self._extract_docx(file_content)
            elif mime_type == "text/plain":
                return self._extract_txt(file_content)
            elif mime_type in ["application/vnd.ms-powerpoint",
                               "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
                return self._extract_pptx(file_content)
            else:
                raise Exception(f"Unsupported MIME type: {mime_type}")
                
        except Exception as e:
            logger.error(f"Failed to extract text from {mime_type}: {str(e)}")
            raise
    
    def _extract_pdf(self, content: bytes) -> str:
        """Extract text from PDF file."""
        try:
            pdf_file = io.BytesIO(content)
            reader = PdfReader(pdf_file)
            
            text_parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            
            full_text = "\n".join(text_parts)
            logger.info(f"Extracted {len(full_text)} characters from PDF")
            return full_text
            
        except Exception as e:
            logger.error(f"PDF extraction failed: {str(e)}")
            raise Exception(f"Failed to extract PDF: {str(e)}")
    
    def _extract_docx(self, content: bytes) -> str:
        """Extract text from DOCX file."""
        try:
            docx_file = io.BytesIO(content)
            doc = Document(docx_file)
            
            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            full_text = "\n".join(text_parts)
            logger.info(f"Extracted {len(full_text)} characters from DOCX")
            return full_text
            
        except Exception as e:
            logger.error(f"DOCX extraction failed: {str(e)}")
            raise Exception(f"Failed to extract DOCX: {str(e)}")
    
    def _extract_txt(self, content: bytes) -> str:
        """Extract text from TXT file."""
        try:
            text = content.decode('utf-8', errors='ignore')
            logger.info(f"Extracted {len(text)} characters from TXT")
            return text
            
        except Exception as e:
            logger.error(f"TXT extraction failed: {str(e)}")
            raise Exception(f"Failed to extract TXT: {str(e)}")
    
    def _extract_pptx(self, content: bytes) -> str:
        """Extract text from PPTX file."""
        try:
            pptx_file = io.BytesIO(content)
            prs = Presentation(pptx_file)
            
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            
            full_text = "\n".join(text_parts)
            logger.info(f"Extracted {len(full_text)} characters from PPTX")
            return full_text
            
        except Exception as e:
            logger.error(f"PPTX extraction failed: {str(e)}")
            raise Exception(f"Failed to extract PPTX: {str(e)}")
    
    def chunk_text(self, text: str, chunk_size: int = 750, overlap: int = 100) -> List[str]:
        """
        Split text into chunks with overlap.
        
        Args:
            text: Text to chunk
            chunk_size: Target chunk size in tokens (approximate)
            overlap: Number of tokens to overlap between chunks
            
        Returns:
            List of text chunks
        """
        if not text or not text.strip():
            return []
        
        # Simple tokenization: split by whitespace
        words = text.split()
        
        if len(words) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(words):
            end = start + chunk_size
            chunk_words = words[start:end]
            
            # Try to break at sentence boundaries
            chunk_text = " ".join(chunk_words)
            
            # If not the last chunk, try to find a sentence boundary
            if end < len(words):
                # Look for sentence endings
                sentences = re.split(r'[.!?]\s+', chunk_text)
                if len(sentences) > 1:
                    # Keep all but the last incomplete sentence
                    chunk_text = ". ".join(sentences[:-1]) + "."
                    # Recalculate actual words used
                    actual_words = len(chunk_text.split())
                    start += actual_words
                else:
                    start = end
            else:
                start = end
            
            chunks.append(chunk_text.strip())
            
            # Move back by overlap amount for next chunk
            if start < len(words):
                start -= overlap
        
        logger.info(f"Split text into {len(chunks)} chunks")
        return chunks
    
    def is_multimodal(self, mime_type: str) -> bool:
        """
        Check if the MIME type is natively supported by Gemini Embedding 2.0.
        """
        multimodal_types = [
            "image/png", "image/jpeg", "image/jpg", "image/webp", "image/heic", "image/heif",
            "video/mp4", "video/mov", "video/quicktime",
            "audio/mpeg", "audio/mp3", "audio/wav", "audio/ogg", "audio/aac",
            "application/pdf"
        ]
        return mime_type in multimodal_types

    def count_tokens(self, text: str) -> int:
        """
        Approximate token count (simple word-based).
        
        Args:
            text: Text to count tokens
            
        Returns:
            Approximate token count
        """
        return len(text.split())


# Global file processor instance
file_processor = FileProcessor()
